from __future__ import annotations

import logging
from pathlib import Path

logger = logging.getLogger(__name__)


def parse_file(path: Path) -> str | None:
    ext = path.suffix.lower()
    try:
        if ext in (".txt", ".md", ".rst"):
            return path.read_text(errors="replace")
        elif ext == ".pdf":
            return _parse_pdf(path)
        elif ext == ".docx":
            return _parse_docx(path)
        elif ext == ".pptx":
            return _parse_pptx(path)
        elif ext == ".xlsx":
            return _parse_xlsx(path)
        else:
            return None
    except Exception:
        logger.exception("Failed to parse %s", path)
        return None


def _parse_pdf(path: Path) -> str:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    parts = []
    for page in reader.pages:
        text = page.extract_text()
        if text:
            parts.append(text)
    return "\n".join(parts)


def _parse_docx(path: Path) -> str:
    import docx

    doc = docx.Document(str(path))
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _parse_pptx(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        parts.append(text)
    return "\n".join(parts)


def _parse_xlsx(path: Path) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(str(path), read_only=True, data_only=True)
    parts = []
    for sheet in wb.worksheets:
        for row in sheet.iter_rows(values_only=True):
            row_text = " ".join(str(cell) for cell in row if cell is not None)
            if row_text.strip():
                parts.append(row_text)
    wb.close()
    return "\n".join(parts)
